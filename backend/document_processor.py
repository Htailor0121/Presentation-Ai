# document_processor.py
import os
import httpx
import fitz  # PyMuPDF
import mammoth
import pandas as pd
import io
import tempfile
from typing import Dict, Any
from bs4 import BeautifulSoup
import mimetypes

# NOTE: Tesseract OCR is optionally used for images / scanned PDFs.
# If you want OCR: pip install pytesseract pillow  and install system tesseract binary.
try:
    import pytesseract
    from PIL import Image
    TESSERACT_AVAILABLE = True
except Exception:
    TESSERACT_AVAILABLE = False

# Supported formats (used by main.py)
supported_formats = {".pdf", ".docx", ".doc", ".txt", ".pptx", ".xlsx", ".csv", ".png", ".jpg", ".jpeg", ".bmp", ".tiff"}

async def fetch_url_content(url: str) -> Dict[str, Any]:
    """Fetch a URL and extract readable text (basic)."""
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            r = await client.get(url)
        if r.status_code != 200:
            return {"success": False, "error": f"HTTP {r.status_code}", "content": ""}

        # Attempt to parse HTML -> text
        soup = BeautifulSoup(r.text, "html.parser")

        # Remove scripts/styles
        for tag in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
            tag.decompose()

        # Simple heuristics: prefer article/body text
        main = soup.find("main") or soup.find("article") or soup.body
        text = main.get_text(separator="\n", strip=True) if main else soup.get_text(separator="\n", strip=True)

        # Short summary metrics
        return {
            "success": True,
            "filename": url,
            "content": text,
            "word_count": len(text.split()),
            "char_count": len(text),
        }
    except Exception as e:
        return {"success": False, "error": str(e), "content": ""}


async def ingest_plain_text(text: str, name: str = "pasted_text") -> Dict[str, Any]:
    text = text or ""
    return {"success": True, "filename": name, "content": text, "word_count": len(text.split()), "char_count": len(text)}


async def process_document(path: str, filename: str) -> Dict[str, Any]:
    """
    Process a document from local path.
    Returns dict: { success, filename, content, word_count, char_count, error (opt) }
    """
    try:
        ext = os.path.splitext(filename)[1].lower()
        file_type, _ = mimetypes.guess_type(path)
        file_type = file_type or ""
        print(f"📄 Detected file type (mimetypes): {file_type}")

        content = ""

        # --- Prefer MIME detection first ---
        if "pdf" in file_type:
            content = extract_text_from_pdf(path)
        elif "wordprocessingml" in file_type:
            content = extract_text_from_docx(path)
        elif "vnd.ms-powerpoint" in file_type or "presentationml" in file_type:
            content = extract_text_from_pptx(path)
        elif "spreadsheetml" in file_type or "excel" in file_type or "csv" in file_type:
            content = extract_text_from_spreadsheet(path)
        elif "text" in file_type:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        elif "image" in file_type:
            content = extract_text_from_image(path)

        # --- Fallback: use extension-based detection ---
        if not content.strip():
            if ext == ".pdf":
                content = extract_text_from_pdf(path)
            elif ext in (".docx", ".doc"):
                content = extract_text_from_docx(path)
            elif ext in (".xlsx", ".xls", ".csv"):
                content = extract_text_from_spreadsheet(path)
            elif ext in (".png", ".jpg", ".jpeg", ".bmp", ".tiff"):
                content = extract_text_from_image(path)
            elif ext == ".txt":
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            elif ext == ".pptx":
                content = extract_text_from_pptx(path)
            else:
                with open(path, "rb") as f:
                    raw = f.read()
                try:
                    content = raw.decode("utf-8", errors="ignore")
                except Exception:
                    content = ""

        return {
            "success": True,
            "filename": filename,
            "content": content,
            "word_count": len(content.split()),
            "char_count": len(content),
        }

    except Exception as e:
        print(f" Error processing {filename}: {e}")
        return {"success": False, "error": str(e), "filename": filename, "content": ""}



def extract_text_from_pdf(path: str) -> str:
    """Extract text from PDF using PyMuPDF (fitz)."""
    doc = fitz.open(path)
    parts = []
    for page in doc:
        text = page.get_text("text")
        if text:
            parts.append(text)
    doc.close()
    return "\n\n".join(parts)


def extract_text_from_docx(path: str) -> str:
    """Extract text from DOCX using mammoth. Falls back if not a valid zip."""
    try:
        with open(path, "rb") as f:
            result = mammoth.extract_raw_text(f)
        return result.value or ""
    except Exception as e:
        print(f" mammoth failed: {e} — falling back to basic read()")
        try:
            # Fallback: attempt plain text extraction
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""



def extract_text_from_spreadsheet(path: str) -> str:
    """Extract CSV/XLSX content into CSV-like representation using pandas."""
    ext = os.path.splitext(path)[1].lower()
    if ext == ".csv":
        df = pd.read_csv(path, dtype=str, encoding="utf-8", errors="ignore")
        return df.to_csv(index=False)
    else:
        # Read all sheets and concat
        xls = pd.read_excel(path, sheet_name=None, dtype=str)
        out = []
        for sheet_name, df in xls.items():
            out.append(f"[Sheet: {sheet_name}]")
            out.append(df.to_csv(index=False))
        return "\n\n".join(out)


def extract_text_from_pptx(path: str) -> str:
    """Extract text from PPTX using python-pptx (optional). Use fallback if not installed."""
    try:
        from pptx import Presentation
    except Exception:
        # If python-pptx not installed, fallback to empty content
        return ""
    prs = Presentation(path)
    out = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        out.append("\n".join(slide_text))
    return "\n\n".join(out)


def extract_text_from_image(path: str) -> str:
    """Use pytesseract (if available) to OCR an image. If tesseract not installed, return empty string."""
    if not TESSERACT_AVAILABLE:
        return ""
    img = Image.open(path)
    text = pytesseract.image_to_string(img)
    return text
