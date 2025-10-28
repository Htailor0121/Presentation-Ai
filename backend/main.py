from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
import uvicorn
import os
import uuid
import tempfile
from datetime import datetime
from ai_service import PresentaionAi
import document_processor
from pydantic import BaseModel
from fastapi.middleware.cors import CORSMiddleware
import asyncio
from fastapi.responses import StreamingResponse
import io
import base64
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import requests
from PIL import Image

class UrlIngestRequest(BaseModel):
    url: str

class TextIngestRequest(BaseModel):
    text: str
    name: str | None = None

class OutlineRequest(BaseModel):
    content: str

class OutlineToSlidesRequest(BaseModel):
    outline: dict


app = FastAPI(
    title="Presentation AI API",
    description="AI-powered presentation generator API",
    version="1.0.0"
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000", "http://localhost:3001", ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Pydantic models
class SlideRequest(BaseModel):
    prompt: Optional[str] = None
    type: Optional[str] = None
    title: Optional[str] = None
    content: Optional[str] = None
    backgroundColor: Optional[str] = "#ffffff"
    textColor: Optional[str] = "#1f2937"
    layout: Optional[str] = "left"
    imageUrl: Optional[str] = None
    chartUrl: Optional[str] = None 


class PresentationRequest(BaseModel):
    title: str
    description: str
    slides: List[SlideRequest]
    theme: str = "modern"

class GeneratePresentationRequest(BaseModel):
    prompt: str
    model: Optional[str] = None
    theme: Optional[str] = None
    include_interactive: bool = True
    num_slides: Optional[int] = None  #  Added for control

class GenerateSlideRequest(BaseModel):
    prompt: str
    model: Optional[str] = None
    theme: Optional[str] = None

class EnhanceSlideRequest(BaseModel):
    slide: dict
    enhancement_type: str = "content"  # content, layout, or overall
    model: Optional[str] = None

class CreateThemeRequest(BaseModel):
    name: str
    primary_color: str
    secondary_color: str
    accent_color: str
    background_color: str
    text_color: str
    font_family: str = "Inter, sans-serif"
    image_style_keywords: List[str] = []

class PresentationResponse(BaseModel):
    title: str
    description: str
    slides: List[SlideRequest]
    theme: str

# AI service instance
ai_service = PresentaionAi()

#  UPDATED: AI service for generating presentations
async def generate_presentation_from_prompt(
    prompt: str, 
    model: str = None, 
    theme: str = None, 
    include_interactive: bool = True,
    num_slides: int = None
) -> PresentationResponse:
    """
    Generate presentation with 8-15 slides.
    If num_slides is None, AI decides (but always 8-15 range).
    """
    try:
        # Default to 10 slides if not specified
        if num_slides is None:
            num_slides = 10
        
        # Ensure 8-15 range
        num_slides = max(8, min(15, num_slides))
        
        print(f"🎯 Generating {num_slides} slides...")
        
        #  Use AI service with num_slides parameter
        ai_response = await ai_service.generate_presentation(
            prompt, 
            model, 
            theme, 
            include_interactive,
            num_slides=num_slides
        )
        
        # Convert to frontend format
        slides = []
        # NEW CODE - USE THIS 
        for slide_data in ai_response.get("slides", []):
            if "id" not in slide_data or not slide_data["id"]:
                slide_data["id"] = f"slide_{uuid.uuid4()}"
    
    #  FIX: Check if slide has a chart
            has_chart = slide_data.get("chartData", {}).get("needed", False)
            image_url = slide_data.get("imageUrl", "")
    
    #  Only generate new image if NO chart AND NO existing image
            if not has_chart and not image_url:
                image_prompt = slide_data.get("imagePrompt")
                if not image_prompt:
                    title = slide_data.get("title", "")
                    content = slide_data.get("content", "")
                    image_prompt = f"{title}. {content[:100]}. Professional, modern, high quality"
        
                try:
                    image_url = await ai_service.generate_image(image_prompt)
                    if not image_url:
                        image_url = "https://source.unsplash.com/1200x800/?presentation,professional"
                except Exception:
                    image_url = "https://source.unsplash.com/1200x800/?presentation,professional"
            elif has_chart:
        #  Chart slides don't need images
                image_url = ""
    
            slide = SlideRequest(
                type=slide_data.get("type", "content"),
                title=slide_data.get("title", ""),
                content=slide_data.get("content", ""),
                backgroundColor=slide_data.get("backgroundColor", "#ffffff"),
                textColor=slide_data.get("textColor", "#1f2937"),
                layout=slide_data.get("layout", "left"),
                imageUrl=image_url,
                chartUrl=slide_data.get("chartUrl", "")
            )
            slides.append(slide)
        
        final_count = len(slides)
        print(f" Created {final_count} slides successfully")
        
        return PresentationResponse(
            title=ai_response.get("title", "Untitled Presentation"),
            description=ai_response.get("description", prompt),
            slides=slides,
            theme=ai_response.get("theme", "modern")
        )
        
    except Exception as e:
        print(f" Error in generate_presentation_from_prompt: {e}")
        # Fallback with 8 slides minimum
        return create_fallback_presentation(prompt, num_slides or 8)
def download_image(url: str) -> str:
    """Download image and return local path"""
    try:
        if url.startswith('data:image'):
            # Handle base64 images
            header, encoded = url.split(',', 1)
            image_data = base64.b64decode(encoded)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                tmp.write(image_data)
                return tmp.name
        else:
            # Download from URL
            response = requests.get(url, timeout=10)
            if response.status_code == 200:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp:
                    tmp.write(response.content)
                    return tmp.name
    except Exception as e:
        print(f"⚠️ Image download error: {e}")
    return None

def hex_to_rgb(hex_color: str):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
#  UPDATED: Fallback with exact slide count
def create_fallback_presentation(prompt: str, num_slides: int = 8) -> PresentationResponse:
    """Create fallback presentation with exact slide count (8-15)"""
    num_slides = max(8, min(15, num_slides))
    title = extract_title_from_prompt(prompt)
    
    slides = []
    
    # Title slide
    slides.append(SlideRequest(
        type="title",
        title=title,
        content=f"A comprehensive presentation about {prompt}",
        backgroundColor="#3b82f6",
        textColor="#ffffff",
        layout="center",
        imageUrl="https://source.unsplash.com/1200x800/?presentation,title"
    ))
    
    # Content slides
    for i in range(1, num_slides - 1):
        slides.append(SlideRequest(
            type="content",
            title=f"Key Point {i}: {title}",
            content=f"Important information about {prompt}\n\n• Detail 1\n• Detail 2\n• Detail 3",
            backgroundColor="#ffffff",
            textColor="#1f2937",
            layout="left",
            imageUrl=f"https://source.unsplash.com/1200x800/?business,presentation,slide{i}"
        ))
    
    # Conclusion slide
    slides.append(SlideRequest(
        type="content",
        title="Conclusion",
        content=f"Thank you!\n\nKey takeaways about {prompt}",
        backgroundColor="#1f2937",
        textColor="#ffffff",
        layout="center",
        imageUrl="https://source.unsplash.com/1200x800/?conclusion,thankyou"
    ))
    
    print(f" Created fallback presentation with {len(slides)} slides")
    
    return PresentationResponse(
        title=title,
        description=prompt,
        slides=slides,
        theme="modern"
    )

def extract_title_from_prompt(prompt: str) -> str:
    """Extract a clean title from the prompt"""
    words = prompt.split()
    title = " ".join(word.capitalize() for word in words)
    title = "".join(c for c in title if c.isalnum() or c.isspace()).strip()
    return title[:50]

# ============================================================================
# API ROUTES
# ============================================================================

@app.get("/")
async def root():
    return {"message": "Presentation AI API", "status": "running"}

@app.get("/api/health")
async def health_check():
    return {
        "status": "OK", 
        "timestamp": datetime.now().isoformat(),
        "service": "Presentation AI API"
    }

#  UPDATED: Generate complete presentation
@app.post("/api/generate-presentation", response_model=PresentationResponse)
async def generate_presentation(request: GeneratePresentationRequest):
    """
    Generate complete presentation (8-15 slides) from prompt.
    Used when user clicks "Generate Slides" in OutlinePage.
    """
    try:
        if not request.prompt.strip():
            raise HTTPException(status_code=400, detail="Prompt is required")
        
        print(f"🎨 Generating full presentation from prompt...")
        
        #  Pass num_slides parameter
        presentation = await generate_presentation_from_prompt(
            request.prompt, 
            request.model, 
            request.theme, 
            request.include_interactive,
            num_slides=request.num_slides  #  Now properly passed
        )
        
        slide_count = len(presentation.slides)
        
        #  Validate slide count
        if slide_count < 8 or slide_count > 15:
            print(f" Warning: Generated {slide_count} slides (expected 8-15)")
        
        print(f" Generated presentation with {slide_count} slides")
        
        return presentation
        
    except Exception as e:
        print(f" Error in generate_presentation: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate presentation: {str(e)}")

@app.post("/api/save-presentation")
async def save_presentation(presentation: PresentationRequest):
    """Save a presentation (placeholder for future database integration)"""
    try:
        return {
            "message": "Presentation saved successfully",
            "id": f"pres_{datetime.now().timestamp()}",
            "timestamp": datetime.now().isoformat()
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save presentation: {str(e)}")

@app.get("/api/presentations")
async def get_presentations():
    """Get all presentations (placeholder for future database integration)"""
    return {
        "presentations": [],
        "message": "No presentations found (database not implemented yet)"
    }

@app.get("/api/models")
async def get_available_models():
    """Get available AI models from OpenRouter"""
    try:
        models = await ai_service.get_available_models()
        return {
            "models": models,
            "default_model": ai_service.default_model
        }
    except Exception as e:
        print(f"Error fetching models: {e}")
        return {
            "models": [],
            "default_model": ai_service.default_model,
            "error": "Failed to fetch models"
        }

@app.post("/api/generate-image")
async def generate_image(request: dict):
    """Generate an image using AI (Hugging Face or Pollinations)"""
    try:
        image_prompt = request.get("prompt")
        if not image_prompt:
            raise HTTPException(status_code=400, detail="Image prompt is required")
        
        image_url = await ai_service.generate_image(image_prompt)
        if not image_url:
            raise HTTPException(status_code=500, detail="Failed to generate image")
        
        return {
            "image_url": image_url,
            "prompt": image_prompt
        }
    except Exception as e:
        print(f"Error generating image: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate image: {str(e)}")

@app.post("/api/upload-document")
async def upload_document(file: UploadFile = File(...)):
    """Upload and process a document (PDF, DOCX, PPT, TXT, CSV, XLSX, etc.)"""
    try:
        # Validate file type
        file_ext = os.path.splitext(file.filename)[1].lower()
        if file_ext not in document_processor.supported_formats:
            raise HTTPException(
                status_code=400, 
                detail=f"Unsupported file format. Supported: {', '.join(document_processor.supported_formats)}"
            )
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            content = await file.read()
            tmp_file.write(content)
            tmp_file_path = tmp_file.name
        
        try:
            # Process the document
            result = await document_processor.process_document(tmp_file_path, file.filename)
            
            if not result["success"]:
                raise HTTPException(status_code=500, detail=result["error"])
            
            return {
                "filename": result["filename"],
                "content": result["content"],
                "word_count": result["word_count"],
                "char_count": result["char_count"]
            }
            
        finally:
            # Clean up temporary file
            if os.path.exists(tmp_file_path):
                os.unlink(tmp_file_path)
                
    except Exception as e:
        print(f"Error processing document: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to process document: {str(e)}")

@app.post("/api/ingest-url")
async def ingest_url(req: UrlIngestRequest):
    """Fetch and extract content from a URL"""
    try:
        if not req.url:
            raise HTTPException(status_code=400, detail="URL is required")
        result = await document_processor.fetch_url_content(req.url)
        if not result.get("success"):
            raise HTTPException(status_code=500, detail=result.get("error", "Failed to fetch URL"))
        return result
    except Exception as e:
        print(f"Error ingesting URL: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to ingest URL: {str(e)}")

@app.post("/api/ingest-text")
async def ingest_text(req: TextIngestRequest):
    """Process plain text input"""
    try:
        if not req.text:
            raise HTTPException(status_code=400, detail="Text is required")
        result = await document_processor.ingest_plain_text(req.text, req.name or "pasted_text")
        if not result.get("success"):
            raise HTTPException(status_code=500, detail=result.get("error", "Failed to ingest text"))
        return result
    except Exception as e:
        print(f"Error ingesting text: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to ingest text: {str(e)}")

#  UPDATED: Generate outline (8-15 sections)
@app.post("/api/generate-outline")
async def generate_outline(req: OutlineRequest):
    """Generate 8-15 section outline from text input"""
    try:
        print(f"🧠 Generating outline from prompt: {req.content[:100]}...")

        if not req.content:
            raise HTTPException(status_code=400, detail="Content is required")

        outline = await ai_service.generate_outline(req.content)

        if not outline or not outline.get("sections"):
            raise HTTPException(status_code=500, detail="Failed to generate outline")

        section_count = len(outline.get("sections", []))
        
        #  Validate count
        if section_count < 8 or section_count > 15:
            print(f" Warning: Got {section_count} sections (expected 8-15)")
        
        print(f" Generated outline with {section_count} sections")

        return outline
        
    except Exception as e:
        print(f" Error generating outline: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate outline: {str(e)}")

#  KEPT: Generate slides from outline
@app.post("/api/generate-slides")
async def generate_slides(req: OutlineToSlidesRequest):
    """Convert outline sections to full slides with images"""
    try:
        if not req.outline:
            raise HTTPException(status_code=400, detail="Outline is required")
        
        print(f"🎨 Converting outline to slides...")
        
        slides_payload = await ai_service.generate_slides_from_outline(req.outline)

        # Ensure images for each slide
        enriched_slides = []
        for idx, slide in enumerate(slides_payload.get("slides", [])):
            image_url = slide.get("imageUrl")
            
            if not image_url:
                prompt = slide.get("imagePrompt") or f"Professional presentation slide: {slide.get('title','')}. {slide.get('content','')[:100]}"
                try:
                    gen = await ai_service.generate_image(prompt)
                    image_url = gen or f"https://source.unsplash.com/1200x800/?presentation,slide{idx}"
                except Exception:
                    image_url = f"https://source.unsplash.com/1200x800/?presentation,slide{idx}"
            
            slide["imageUrl"] = image_url
            enriched_slides.append(slide)

        slides_payload["slides"] = enriched_slides
        
        slide_count = len(enriched_slides)
        print(f" Generated {slide_count} complete slides with images")
        
        return slides_payload
        
    except Exception as e:
        print(f" Error generating slides from outline: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate slides: {str(e)}")

#  UPDATED: Summarize document (Critical endpoint)
@app.post("/api/summarize-document")
async def summarize_document(request: dict):
    """
    Create outline or full slides from document content.
    
    outline_only=True  → Returns 8-15 section outline (fast, no images)
    outline_only=False → Returns 8-15 complete slides with images (slower)
    """
    try:
        document_content = request.get("content")
        filename = request.get("filename", "document")
        outline_only = request.get("outline_only", False)

        if not document_content:
            raise HTTPException(status_code=400, detail="Document content is required")

        print(f"📄 Processing document: {filename}")
        print(f"   Mode: {'Outline only (8-15 sections)' if outline_only else 'Full slides with images (8-15)'}")

        # Get AI response
        ai_response = await ai_service.summarize_document(
            document_content, 
            filename, 
            outline_only
        )

        # ============================================================
        # MODE A: Outline Only (Fast)
        # ============================================================
        if outline_only:
            sections = ai_response.get("sections", [])
            section_count = len(sections)
            
            #  Validate count
            if section_count < 8:
                print(f" Warning: Only {section_count} sections, expected min 8")
            elif section_count > 15:
                print(f" Warning: Got {section_count} sections, expected max 15")

            print(f" Returning {section_count} outline sections")

            return {
                "title": ai_response.get("title", f"Outline: {filename}"),
                "description": ai_response.get("description", "Document outline"),
                "outline": sections,  # Frontend expects 'outline' key
                "sections": sections   # Also 'sections' for compatibility
            }

        # ============================================================
        # MODE B: Full Slides with Images (Slower)
        # ============================================================
        else:
            slides = ai_response.get("slides", [])
            
            if not slides:
                raise HTTPException(status_code=500, detail="No slides generated from document")

            slide_count = len(slides)
            
            #  Validate count
            if slide_count < 8:
                print(f" Warning: Only {slide_count} slides, expected min 8")
            elif slide_count > 15:
                print(f" Warning: Got {slide_count} slides, expected max 15")

            print(f"🎨 Generating images for {slide_count} slides...")

            # Add images to each slide
            enriched_slides = []
            for idx, slide_data in enumerate(slides):
                print(f"  📍 [{idx+1}/{slide_count}] Processing: {slide_data.get('title', 'Untitled')[:50]}")
                
                # Get or generate image
                image_url = slide_data.get("imageUrl")
                
                if not image_url:
                    title = slide_data.get("title", "")
                    content = slide_data.get("content", "")[:200]
                    
                    image_prompt = f"Professional presentation slide for: {title}. Modern, clean, high quality, 4k"
                    
                    try:
                        image_url = await ai_service.generate_image(image_prompt)
                        if not image_url:
                            image_url = f"https://source.unsplash.com/1200x800/?{title.replace(' ', ',')},professional"
                        print(f"     Image generated")
                    except Exception as img_error:
                        print(f"     Image error: {img_error}")
                        image_url = f"https://source.unsplash.com/1200x800/?presentation,slide{idx}"
                
                # Build clean slide object
                enriched_slide = {
                    "id": slide_data.get("id", f"slide_{idx}_{int(datetime.now().timestamp())}"),
                    "type": slide_data.get("type", "content"),
                    "title": slide_data.get("title", f"Slide {idx+1}"),
                    "content": slide_data.get("content", ""),
                    "imageUrl": image_url,
                    "chartUrl": slide_data.get("chartUrl", ""),
                    "layout": slide_data.get("layout", "split"),
                    "textAlign": slide_data.get("textAlign", "left"),
                    "backgroundColor": slide_data.get("backgroundColor", "#ffffff"),
                    "textColor": slide_data.get("textColor", "#1f2937")
                }
                
                enriched_slides.append(enriched_slide)

            print(f" Generated {len(enriched_slides)} complete slides with images")

            return {
                "title": ai_response.get("title", f"Summary: {filename}"),
                "description": ai_response.get("description", f"Presentation from {filename}"),
                "slides": enriched_slides,
                "theme": ai_response.get("theme", "modern")
            }

    except HTTPException:
        raise
    except Exception as e:
        print(f" Error in summarize_document: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Failed to process document: {str(e)}")

# ============================================================================
# SLIDE ENHANCEMENT ENDPOINTS
# ============================================================================

@app.post("/api/enhance-slide")
async def enhance_slide(data: SlideRequest):
    """Enhance slide content with AI"""
    try:
        result = await ai_service.generate_ai_text(
            f"Enhance and expand this slide content professionally:\n\n{data.prompt}"
        )
        return {"enhanced": result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/rewrite-slide")
async def rewrite_slide(data: SlideRequest):
    """Rewrite slide with better flow"""
    try:
        result = await ai_service.generate_ai_text(
            f"Rewrite this slide with better flow and clarity:\n\n{data.prompt}"
        )
        return {"rewritten": result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/summarize-slide")
async def summarize_slide(data: SlideRequest):
    """Summarize slide content"""
    try:
        result = await ai_service.generate_ai_text(
            f"Summarize this slide in a concise and clear way:\n\n{data.prompt}"
        )
        
        # Clean out model instruction tags
        cleaned = (
            result.replace("<s>", "")
                  .replace("</s>", "")
                  .replace("[B_INST]", "")
                  .replace("[/B_INST]", "")
                  .strip()
        )
        
        return {"summary": cleaned}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/expand-slide")
async def expand_slide(data: SlideRequest):
    """Expand slide with more details"""
    try:
        result = await ai_service.generate_ai_text(
            f"Expand this slide by adding more details and examples:\n\n{data.prompt}"
        )
        return {"expanded": result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/change-tone")
async def change_tone(data: SlideRequest):
    """Change slide tone"""
    try:
        result = await ai_service.generate_ai_text(
            f"Rewrite this slide with a more friendly and engaging tone:\n\n{data.prompt}"
        )
        return {"tone_changed": result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ============================================================================
# THEME ENDPOINTS
# ============================================================================

@app.get("/api/themes")
async def get_themes():
    """Get available themes"""
    try:
        themes = ai_service.get_available_themes()
        return {
            "themes": themes,
            "message": "Available themes retrieved successfully"
        }
    except Exception as e:
        print(f"Error fetching themes: {e}")
        return {
            "themes": [],
            "error": "Failed to fetch themes"
        }

@app.post("/api/create-theme")
async def create_theme(request: CreateThemeRequest):
    """Create a custom theme"""
    try:
        if not request.name.strip():
            raise HTTPException(status_code=400, detail="Theme name is required")
        
        theme_name = ai_service.create_custom_theme(
            request.name,
            primary_color=request.primary_color,
            secondary_color=request.secondary_color,
            accent_color=request.accent_color,
            background_color=request.background_color,
            text_color=request.text_color,
            font_family=request.font_family,
            image_style_keywords=request.image_style_keywords
        )
        
        return {
            "theme_name": theme_name,
            "message": "Custom theme created successfully"
        }
        
    except Exception as e:
        print(f"Error creating theme: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to create theme: {str(e)}")

# ============================================================================
# SINGLE SLIDE GENERATION (Gamma-style)
# ============================================================================

@app.post("/api/generate-slide")
async def generate_slide(request: GenerateSlideRequest):
    """Generate a single slide from a prompt (Gamma-style)"""
    try:
        if not request.prompt.strip():
            raise HTTPException(status_code=400, detail="Prompt is required")
        
        slide = await ai_service.generate_slide_from_prompt(
            request.prompt, 
            request.model, 
            request.theme
        )
        return slide
        
    except Exception as e:
        print(f"Error in generate_slide endpoint: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate slide: {str(e)}")

@app.post("/api/export-presentation")
async def export_presentation_endpoint(request: dict):
    """Export presentation in multiple formats (JSON, PDF, PPTX)"""
    try:
        title = request.get("title", "Untitled Presentation")
        slides = request.get("slides", [])
        theme = request.get("theme", "modern")
        export_format = request.get("format", "json").lower()
        
        # Sanitize filename
        safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_title = safe_title.replace(' ', '-')[:50] or "presentation"
        
        if export_format == "pdf":
            return await export_as_pdf(title, slides, safe_title)
        elif export_format == "pptx":
            return await export_as_pptx(title, slides, safe_title)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {export_format}")
            
    except Exception as e:
        print(f"❌ Export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

async def export_as_pdf(title: str, slides: list, filename: str):
    """Export as PDF file"""
    try:
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            pdf_path = tmp_file.name
        
        # Create PDF
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter
        
        for idx, slide in enumerate(slides):
            if idx > 0:
                c.showPage()
            
            # Background color
            bg_color = slide.get('backgroundColor', '#ffffff')
            if bg_color and bg_color != '#ffffff':
                try:
                    r, g, b = hex_to_rgb(bg_color)
                    c.setFillColorRGB(r/255, g/255, b/255)
                    c.rect(0, 0, width, height, fill=1, stroke=0)
                except:
                    pass
            
            # Title
            c.setFont("Helvetica-Bold", 28)
            text_color = slide.get('textColor', '#1f2937')
            try:
                r, g, b = hex_to_rgb(text_color)
                c.setFillColorRGB(r/255, g/255, b/255)
            except:
                c.setFillColorRGB(0.12, 0.16, 0.22)
            
            title_text = slide.get('title', f'Slide {idx + 1}')
            c.drawString(50, height - 80, title_text[:60])
            
            # Content
            c.setFont("Helvetica", 14)
            y_position = height - 130
            content = slide.get('content', '')
            content_lines = content.split('\n')
            
            for line in content_lines[:20]:
                if y_position < 150:
                    break
                clean_line = line.replace('•', '-').strip()[:90]
                if clean_line:
                    c.drawString(70, y_position, clean_line)
                y_position -= 22
            
            # Image (if layout includes image and not just chart)
            layout = slide.get('layout', 'split')
            image_url = slide.get('imageUrl')
            chart_url = slide.get('chartUrl')
            
            if image_url and layout != 'full-text' and not chart_url:
                img_path = download_image(image_url)
                if img_path:
                    try:
                        img = ImageReader(img_path)
                        img_width = 250
                        img_height = 180
                        
                        if layout == 'split':
                            c.drawImage(img, width - img_width - 50, height - 280, 
                                      width=img_width, height=img_height, 
                                      preserveAspectRatio=True, mask='auto')
                        else:
                            c.drawImage(img, (width - img_width) / 2, 200, 
                                      width=img_width, height=img_height, 
                                      preserveAspectRatio=True, mask='auto')
                        
                        os.unlink(img_path)
                    except Exception as e:
                        print(f"⚠️ PDF image error: {e}")
            
            # Chart (takes priority over image)
            if chart_url:
                chart_path = download_image(chart_url)
                if chart_path:
                    try:
                        chart_img = ImageReader(chart_path)
                        chart_width = 500
                        chart_height = 300
                        
                        c.drawImage(chart_img, (width - chart_width) / 2, 100,
                                  width=chart_width, height=chart_height,
                                  preserveAspectRatio=True, mask='auto')
                        
                        os.unlink(chart_path)
                    except Exception as e:
                        print(f"⚠️ PDF chart error: {e}")
            
            # Footer
            c.setFont("Helvetica", 9)
            c.setFillColorRGB(0.6, 0.6, 0.6)
            c.drawString(50, 30, f"{title} - Slide {idx + 1}/{len(slides)}")
            c.drawRightString(width - 50, 30, datetime.now().strftime("%B %d, %Y"))
        
        c.save()
        
        # Read and return file
        with open(pdf_path, 'rb') as f:
            pdf_data = f.read()
        
        os.unlink(pdf_path)
        
        return StreamingResponse(
            io.BytesIO(pdf_data),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}.pdf"'
            }
        )
        
    except Exception as e:
        print(f"❌ PDF export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PDF export failed: {str(e)}")

async def export_as_pptx(title: str, slides: list, filename: str):
    """Export as PowerPoint file"""
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for idx, slide_data in enumerate(slides):
            # Use blank layout
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Background color
            background = slide.background
            fill = background.fill
            fill.solid()
            bg_color = slide_data.get('backgroundColor', '#ffffff')
            try:
                r, g, b = hex_to_rgb(bg_color)
                fill.fore_color.rgb = RGBColor(r, g, b)
            except:
                fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            layout = slide_data.get('layout', 'split')
            
            # Title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            p = title_frame.paragraphs[0]
            p.text = slide_data.get('title', f'Slide {idx + 1}')
            p.font.size = Pt(36)
            p.font.bold = True
            
            try:
                text_color = slide_data.get('textColor', '#1f2937')
                r, g, b = hex_to_rgb(text_color)
                p.font.color.rgb = RGBColor(r, g, b)
            except:
                p.font.color.rgb = RGBColor(31, 41, 55)
            
            text_align = slide_data.get('textAlign', 'left')
            if text_align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif text_align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            
            # Content positioning based on layout
            if layout == 'centered' or layout == 'center':
                content_left = Inches(1.5)
                content_top = Inches(2.5)
                content_width = Inches(7)
                content_height = Inches(4)
            elif layout == 'full-text':
                content_left = Inches(0.5)
                content_top = Inches(2)
                content_width = Inches(9)
                content_height = Inches(5)
            else:  # split or default
                content_left = Inches(0.5)
                content_top = Inches(2)
                content_width = Inches(4.5)
                content_height = Inches(5)
            
            # Content
            content_box = slide.shapes.add_textbox(
                content_left, content_top, content_width, content_height
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            content_text = slide_data.get('content', '')
            content_lines = content_text.split('\n')
            
            for line_idx, line in enumerate(content_lines[:15]):
                if line_idx > 0:
                    content_frame.add_paragraph()
                p = content_frame.paragraphs[line_idx]
                p.text = line.strip()
                p.font.size = Pt(18)
                p.space_after = Pt(8)
                
                try:
                    r, g, b = hex_to_rgb(text_color)
                    p.font.color.rgb = RGBColor(r, g, b)
                except:
                    p.font.color.rgb = RGBColor(55, 65, 81)
                
                if text_align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif text_align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT
            
            # Chart (priority)
            chart_url = slide_data.get('chartUrl')
            if chart_url:
                chart_path = download_image(chart_url)
                if chart_path:
                    try:
                        slide.shapes.add_picture(
                            chart_path,
                            Inches(1),
                            Inches(2.5),
                            width=Inches(8),
                            height=Inches(4.5)
                        )
                        os.unlink(chart_path)
                    except Exception as e:
                        print(f"⚠️ PPTX chart error: {e}")
            
            # Image (only if no chart and layout supports it)
            elif layout not in ['full-text', 'centered', 'center']:
                image_url = slide_data.get('imageUrl')
                if image_url:
                    img_path = download_image(image_url)
                    if img_path:
                        try:
                            slide.shapes.add_picture(
                                img_path,
                                Inches(5.5),
                                Inches(2),
                                width=Inches(4),
                                height=Inches(5)
                            )
                            os.unlink(img_path)
                        except Exception as e:
                            print(f"⚠️ PPTX image error: {e}")
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            pptx_path = tmp_file.name
        
        prs.save(pptx_path)
        
        # Read and return
        with open(pptx_path, 'rb') as f:
            pptx_data = f.read()
        
        os.unlink(pptx_path)
        
        return StreamingResponse(
            io.BytesIO(pptx_data),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}.pptx"'
            }
        )
        
    except Exception as e:
        print(f"❌ PPTX export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"PPTX export failed: {str(e)}")

# ============================================================================
# SERVE STATIC FILES (Production)
# ============================================================================

if os.path.exists("dist"):
    app.mount("/static", StaticFiles(directory="dist"), name="static")
    
    @app.get("/{full_path:path}")
    async def serve_spa(full_path: str):
        """Serve the React app for any non-API routes"""
        if full_path.startswith("api/"):
            return {"error": "API endpoint not found"}
        return FileResponse("dist/index.html")

# ============================================================================
# RUN SERVER
# ============================================================================

if __name__ == "__main__":
    uvicorn.run(
        "main:app", 
        host="0.0.0.0", 
        port=5000, 
        reload=True,
        log_level="info"
    )